"""Unified text extraction across PDF / DOCX / XLSX / PPTX / images / md / txt.

`extract_text(path)` returns plain UTF-8 text suitable for stuffing into an
LLM prompt. PDFs without a usable text layer fall back to OCR. Anything we
don't recognise raises `UnsupportedFileType`.
"""

from __future__ import annotations

from pathlib import Path

# Threshold below which a PDF is considered "scanned" and we fall back to OCR.
_PDF_TEXT_FLOOR = 50


class UnsupportedFileType(Exception):
    pass


def extract_text(path: str | Path) -> str:
    p = Path(path)
    suffix = p.suffix.lower()

    if suffix == ".pdf":
        return _extract_pdf(p)
    if suffix == ".docx":
        return _extract_docx(p)
    if suffix in {".xlsx", ".xls"}:
        return _extract_xlsx(p)
    if suffix == ".pptx":
        return _extract_pptx(p)
    if suffix in {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp"}:
        return _extract_image(p)
    if suffix in {".md", ".txt"}:
        return p.read_text(encoding="utf-8", errors="replace")

    raise UnsupportedFileType(f"unsupported file type: {suffix or '(no extension)'}")


# --------------------------------------------------------------------------- #
# Per-format extractors
# --------------------------------------------------------------------------- #

def _extract_pdf(p: Path) -> str:
    import pdfplumber

    parts: list[str] = []
    with pdfplumber.open(str(p)) as pdf:
        for page in pdf.pages:
            t = page.extract_text() or ""
            if t:
                parts.append(t)
    text = "\n".join(parts).strip()
    if len(text) >= _PDF_TEXT_FLOOR:
        return text
    # Probably a scan: OCR each page.
    return _ocr_pdf(p)


def _ocr_pdf(p: Path) -> str:
    from pdf2image import convert_from_path
    import pytesseract

    images = convert_from_path(str(p))
    return "\n".join(pytesseract.image_to_string(img) for img in images).strip()


def _extract_docx(p: Path) -> str:
    from docx import Document

    doc = Document(str(p))
    parts: list[str] = [para.text for para in doc.paragraphs if para.text]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append("\t".join(cells))
    
    # Fallback: extract text from VML textboxes if no text found via python-docx
    if not parts:
        vml_texts = _extract_vml_textboxes(p)
        parts.extend(vml_texts)
    
    return "\n".join(parts).strip()


def _extract_vml_textboxes(p: Path) -> list[str]:
    """Extract text from VML textboxes in DOCX by parsing raw XML.
    
    Some DOCX files (especially forms) use VML textboxes (<v:textbox>) to
    contain text instead of standard paragraphs. python-docx does not read
    these, so we parse the raw XML directly.
    """
    import zipfile
    from xml.etree import ElementTree as ET
    
    parts: list[str] = []
    
    # Namespace for WordprocessingML text elements
    W_NS = '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}'
    V_NS = '{urn:schemas-microsoft-com:vml}'
    
    try:
        with zipfile.ZipFile(str(p), 'r') as zf:
            # Read document.xml which contains the main content
            if 'word/document.xml' not in zf.namelist():
                return parts
            
            xml_content = zf.read('word/document.xml')
            root = ET.fromstring(xml_content)
            
            # Find all VML textbox elements and extract text
            for textbox in root.iter(f'{V_NS}textbox'):
                textbox_parts: list[str] = []
                # Text is stored in <w:t> elements inside the textbox
                for t_elem in textbox.iter(f'{W_NS}t'):
                    if t_elem.text:
                        textbox_parts.append(t_elem.text)
                if textbox_parts:
                    parts.append(''.join(textbox_parts))
            
            # Also check for text in <w:txbxContent> (textbox content) elements
            # which may not be inside VML namespace
            for txbx_content in root.iter(f'{W_NS}txbxContent'):
                for t_elem in txbx_content.iter(f'{W_NS}t'):
                    if t_elem.text and t_elem.text not in parts:
                        parts.append(t_elem.text)
                        
    except (zipfile.BadZipFile, ET.ParseError):
        # If we can't parse the file, return empty list
        pass
    
    return parts


def _extract_xlsx(p: Path) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(str(p), data_only=True, read_only=True)
    parts: list[str] = []
    for sheet in wb.worksheets:
        parts.append(f"# Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            cells = ["" if v is None else str(v) for v in row]
            if any(cells):
                parts.append("\t".join(cells))
    return "\n".join(parts).strip()


def _extract_pptx(p: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(p))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts.append(f"# Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = "".join(run.text for run in para.runs).strip()
                    if txt:
                        parts.append(txt)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        parts.append("\t".join(cells))
    return "\n".join(parts).strip()


def _extract_image(p: Path) -> str:
    from PIL import Image
    import pytesseract

    return pytesseract.image_to_string(Image.open(str(p))).strip()
